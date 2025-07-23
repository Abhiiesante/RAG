import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import pandas as pd
from typing import List, Dict, Any
import io

class DocumentParser:
    """Handles parsing of various document formats"""
    
    @staticmethod
    def parse_pdf(file_content: bytes) -> List[Dict[str, Any]]:
        """Parse PDF and extract text chunks"""
        chunks = []
        doc = fitz.open(stream=file_content, filetype="pdf")
        
        for page_num in range(doc.page_count):
            page = doc[page_num]
            text = page.get_text()
            
            if text.strip():
                chunks.append({
                    "content": text,
                    "metadata": {
                        "page": page_num + 1,
                        "type": "pdf"
                    }
                })
        
        doc.close()
        return chunks
    
    @staticmethod
    def parse_pptx(file_content: bytes) -> List[Dict[str, Any]]:
        """Parse PPTX and extract text from slides"""
        chunks = []
        prs = Presentation(io.BytesIO(file_content))
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            if slide_text:
                chunks.append({
                    "content": "\n".join(slide_text),
                    "metadata": {
                        "slide": slide_num + 1,
                        "type": "pptx"
                    }
                })
        
        return chunks
    
    @staticmethod
    def parse_docx(file_content: bytes) -> List[Dict[str, Any]]:
        """Parse DOCX and extract paragraphs"""
        chunks = []
        doc = Document(io.BytesIO(file_content))
        
        for para_num, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                chunks.append({
                    "content": paragraph.text,
                    "metadata": {
                        "paragraph": para_num + 1,
                        "type": "docx"
                    }
                })
        
        return chunks
    
    @staticmethod
    def parse_csv(file_content: bytes) -> List[Dict[str, Any]]:
        """Parse CSV and convert to text chunks"""
        chunks = []
        df = pd.read_csv(io.BytesIO(file_content))
        
        # Convert headers to text
        headers_text = f"CSV Headers: {', '.join(df.columns.tolist())}"
        chunks.append({
            "content": headers_text,
            "metadata": {
                "row": "headers",
                "type": "csv"
            }
        })
        
        # Convert rows to text
        for idx, row in df.iterrows():
            row_text = []
            for col, value in row.items():
                if pd.notna(value):
                    row_text.append(f"{col}: {value}")
            
            if row_text:
                chunks.append({
                    "content": " | ".join(row_text),
                    "metadata": {
                        "row": idx + 1,
                        "type": "csv"
                    }
                })
        
        return chunks
    
    @staticmethod
    def parse_txt(file_content: bytes) -> List[Dict[str, Any]]:
        """Parse TXT/Markdown files"""
        text = file_content.decode('utf-8')
        
        # Split into paragraphs
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        chunks = []
        for i, paragraph in enumerate(paragraphs):
            chunks.append({
                "content": paragraph,
                "metadata": {
                    "paragraph": i + 1,
                    "type": "txt"
                }
            })
        
        return chunks

class DocumentChunker:
    """Handles chunking of document text for better retrieval"""
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings near the chunk boundary
                for i in range(min(100, chunk_size // 4)):
                    if end - i < len(text) and text[end - i] in '.!?\n':
                        end = end - i + 1
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap
            if start >= len(text):
                break
        
        return chunks
